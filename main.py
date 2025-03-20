import os
from openai import OpenAI
from dotenv import load_dotenv
import streamlit as st
import pdfplumber
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
from pptx import Presentation
from docx import Document

# Load API Key
load_dotenv()
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")

# Initialize OpenAI Client for OpenRouter
client = OpenAI(
    base_url="https://openrouter.ai/api/v1",
    api_key=OPENROUTER_API_KEY,
)

# Set page layout
st.set_page_config(page_title="StudyBuddy", layout="wide")

# Sidebar UI
st.sidebar.title("⚙️ Settings")
st.sidebar.info("Upload your study materials and let StudyBuddy AI help you summarize, generate quizzes, and more!")

# Session State Initialization (Persists Data Across Tabs)
if "summary" not in st.session_state:
    st.session_state.summary = ""
if "section_summary" not in st.session_state:
    st.session_state.section_summary = ""
if "quiz" not in st.session_state:
    st.session_state.quiz = ""
if "flashcards" not in st.session_state:
    st.session_state.flashcards = ""
if "concept_map" not in st.session_state:
    st.session_state.concept_map = ""

# Cache OCR results
@st.cache_data
def extract_text_from_pdf(pdf_path):
    text = ""
    with pdfplumber.open(pdf_path) as pdf:
        for page in pdf.pages:
            extracted_text = page.extract_text()
            if extracted_text:
                text += extracted_text + "\n"
    
    if not text.strip():
        st.warning("No text found using pdfplumber, switching to OCR...")
        images = convert_from_path(pdf_path)
        for i, img in enumerate(images):
            st.info(f"Running OCR on page {i+1}...")
            ocr_text = pytesseract.image_to_string(img)
            text += ocr_text + "\n"

    return text.strip() if text.strip() else "Text extraction failed. Try another document."

# Extract text from PPTX
def extract_text_from_pptx(pptx_path):
    presentation = Presentation(pptx_path)
    text = ""
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()

# Extract text from DOCX
def extract_text_from_docx(docx_path):
    doc = Document(docx_path)
    text = "\n".join([para.text for para in doc.paragraphs])
    return text.strip()

# AI Query Function
def ai_query(prompt):
    completion = client.chat.completions.create(
        extra_headers={
            "HTTP-Referer": "https://thestudybuddy.streamlit.app/",  # Optional: replace with your site URL
            "X-Title": "StudyBuddy AI",  # Optional: replace with your app name
        },
        extra_body={},
        model="deepseek/deepseek-r1:free",
        messages=[{"role": "user", "content": prompt}],
    )
    return completion.choices[0].message.content

# AI Functions
def summarize_notes(notes):
    return ai_query(f"Summarize the following study notes effectively and detailed, do not be restricted, be as informative as you can be:\n{notes}")

def summarize_sections(notes):
    return ai_query(f"Break down these notes by sections and summarize each separately:\n{notes}")

def generate_quiz(notes):
    return ai_query(f"Generate multiple-choice quiz questions based on the following notes, do not bold the answer when giving the options, and don't make the answer always option B, make the answers a varied so one is b then another c than a, etc:\n{notes}")

def generate_flashcards(notes):
    return ai_query(f"Generate as many flashcards as you think necessary in Q/A format based on these notes:\n{notes}")

def generate_concept_map(notes):
    return ai_query(f"Create a concept map with key topics and relationships based on:\n{notes}")

# Main UI
st.title("StudyBuddy")
st.write("Upload your study materials and let StudyBuddy generate summaries, quizzes, and flashcards!")

uploaded_files = st.file_uploader(
    "📂 Upload Study Materials (Multiple PDFs, PPTX, DOCX)", 
    type=["pdf", "pptx", "docx"], 
    accept_multiple_files=True
)

# File Processing
if uploaded_files:
    all_notes = ""

    with st.spinner("Processing files..."):
        for uploaded_file in uploaded_files:
            file_extension = uploaded_file.name.split(".")[-1]
            file_path = f"temp_{uploaded_file.name}"

            with open(file_path, "wb") as f:
                f.write(uploaded_file.getbuffer())

            if file_extension == "pdf":
                extracted_text = extract_text_from_pdf(file_path)
            elif file_extension == "pptx":
                extracted_text = extract_text_from_pptx(file_path)
            elif file_extension == "docx":
                extracted_text = extract_text_from_docx(file_path)
            else:
                extracted_text = ""

            all_notes += f"\n\n--- Extracted from {uploaded_file.name} ---\n{extracted_text}\n"

    st.success("✅ Study Materials Processed!")
    st.write("## 📜 Extracted Study Notes")
    st.text_area("Study Notes:", all_notes, height=200)

    # Tabs for better UI navigation
    tab1, tab2, tab3, tab4 = st.tabs(["📖 Summarization", "📝 Quiz", "🎴 Flashcards", "🌍 Concept Map"])

    with tab1:
        col1, col2 = st.columns(2)
        with col1:
            if st.button("📝 Summarize Chapter", use_container_width=True):
                st.session_state.summary = summarize_notes(all_notes)
        with col2:
            if st.button("🗂 Summarize by Sections", use_container_width=True):
                st.session_state.section_summary = summarize_sections(all_notes)

        # Display persistent summary results
        if st.session_state.summary:
            st.write("## 📖 Summary")
            st.write(st.session_state.summary)
        if st.session_state.section_summary:
            st.write("## 📚 Section-wise Summary")
            st.write(st.session_state.section_summary)

    with tab2:
        if st.button("📝 Generate Quiz", use_container_width=True):
            st.session_state.quiz = generate_quiz(all_notes)

        # Display persistent quiz results
        if st.session_state.quiz:
            st.write("## 🎯 Quiz")
            st.write(st.session_state.quiz)

    with tab3:
        if st.button("🎴 Generate Flashcards", use_container_width=True):
            st.session_state.flashcards = generate_flashcards(all_notes)

        # Display persistent flashcards results
        if st.session_state.flashcards:
            st.write("## 🎴 Flashcards")
            st.write(st.session_state.flashcards)

    with tab4:
        if st.button("🌍 Generate Concept Map", use_container_width=True):
            st.session_state.concept_map = generate_concept_map(all_notes)

        # Display persistent concept map results
        if st.session_state.concept_map:
            st.write("## 🌍 Concept Map")
            st.write(st.session_state.concept_map)

    # Q&A Section
    st.write("## ❓ Ask StudyBuddy AI a Question")
    user_question = st.text_input("Ask a question about your study materials")
    if st.button("🤖 Get AI Answer", use_container_width=True):
        answer = ai_query(user_question)
        st.write("## 💡 AI Answer")
        st.write(answer)
